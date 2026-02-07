"""Generate class-facing PowerPoint decks for the Aviation merit badge sessions.

Outputs:
- slides/session-1.pptx
- slides/session-2.pptx
- slides/session-2-req-5.pptx

This script intentionally uses concise, class-facing bullets and puts extra detail
into speaker notes for the instructor.

Usage:
  .venv/bin/python tools/generate_slides.py
"""

from __future__ import annotations

from dataclasses import dataclass
from pathlib import Path
import re
from typing import Iterable, List, Optional

from pptx import Presentation
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
SLIDES_DIR = ROOT / "slides"
INSTRUCTOR_NOTES_DIR = ROOT / "instructor-notes"
SESSION_2_REQ_5_SOURCE = INSTRUCTOR_NOTES_DIR / "instructor-guide-session-2-requirement-5.md"


@dataclass(frozen=True)
class SlideSpec:
    title: str
    bullets: List[str]
    notes: Optional[str] = None


def _add_title_and_bullets(prs: Presentation, spec: SlideSpec) -> None:
    layout = prs.slide_layouts[1]  # Title and Content
    slide = prs.slides.add_slide(layout)

    slide.shapes.title.text = spec.title

    body = slide.shapes.placeholders[1].text_frame
    body.clear()
    for i, bullet in enumerate(spec.bullets):
        p = body.paragraphs[0] if i == 0 else body.add_paragraph()
        p.text = bullet
        p.level = 0

    # Typography tweaks for readability
    for p in body.paragraphs:
        for run in p.runs:
            run.font.size = Pt(28)

    if spec.notes:
        notes = slide.notes_slide.notes_text_frame
        notes.clear()
        notes.text = spec.notes


def _add_title_slide(prs: Presentation, title: str, subtitle: str) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[0])  # Title
    slide.shapes.title.text = title
    slide.placeholders[1].text = subtitle

    # Slightly larger title for room visibility
    for run in slide.shapes.title.text_frame.paragraphs[0].runs:
        run.font.size = Pt(44)


def _add_quick_activity_slide(prs: Presentation, title: str, steps: Iterable[str], notes: str) -> None:
    _add_title_and_bullets(
        prs,
        SlideSpec(
            title=title,
            bullets=[f"Step {i+1}: {s}" for i, s in enumerate(steps)],
            notes=notes,
        ),
    )


def _read_text(path: Path) -> str:
    return path.read_text(encoding="utf-8")


def _extract_markdown_section(markdown: str, heading_text: str) -> str:
    """Return the markdown content under the first matching heading.

    Matches headings like "## Outcomes" or "### 0:03–0:15 – Pilot certificates".
    The returned content ends at the next heading of the same or higher level.
    """

    heading_pattern = re.compile(r"^(#{1,6})\\s+(.*)\\s*$")
    lines = markdown.splitlines()

    target_idx: Optional[int] = None
    target_level: Optional[int] = None

    for idx, line in enumerate(lines):
        match = heading_pattern.match(line)
        if not match:
            continue
        level = len(match.group(1))
        title = match.group(2).strip()
        if title == heading_text.strip():
            target_idx = idx
            target_level = level
            break

    if target_idx is None or target_level is None:
        return ""

    start = target_idx + 1
    end = len(lines)
    for idx in range(start, len(lines)):
        match = heading_pattern.match(lines[idx])
        if not match:
            continue
        level = len(match.group(1))
        if level <= target_level:
            end = idx
            break

    section = "\n".join(lines[start:end]).strip()
    return section


def _notes_from_sections(markdown: str, headings: List[str]) -> str:
    parts: List[str] = []
    for heading in headings:
        content = _extract_markdown_section(markdown, heading)
        if content:
            parts.append(f"[{heading}]\n{content}")
    return "\n\n".join(parts).strip() or None


def build_session_1() -> Presentation:
    prs = Presentation()

    _add_title_slide(
        prs,
        title="Aviation Merit Badge – Session 1",
        subtitle="Aviation basics • How flight works • What pilots see in the cockpit",
    )

    slides: List[SlideSpec] = [
        SlideSpec(
            title="Today’s goals",
            bullets=[
                "What counts as an aircraft",
                "Fixed wing vs rotary wing",
                "Engines: piston vs turbine vs jet",
                "Four forces of flight",
                "Lift + airfoils (Bernoulli + downwash)",
                "Control surfaces (pitch/roll/yaw)",
                "Cockpit instruments (what they tell the pilot)",
            ],
            notes=(
                "Keep this fast. Mention that Session 2 is hands-on with model/drone flight.\n"
                "If running short on time later: compress history and keep time for forces/lift/controls/instruments."
            ),
        ),
        SlideSpec(
            title="What is an aircraft?",
            bullets=[
                "A machine that flies by being supported by air",
                "It can create lift with wings…",
                "…or lift/thrust with rotors",
            ],
            notes="Ask: name 3 aircraft you’ve seen and one thing it’s used for.",
        ),
        SlideSpec(
            title="Kinds of aircraft (examples)",
            bullets=[
                "Fixed-wing airplane (training, transport, cargo)",
                "Rotary-wing (helicopter) (EMS, rescue, lifting)",
                "Other: glider • balloon • jet airliner • drone/UAS",
            ],
            notes="Have scouts categorize examples quickly. Keep it concrete.",
        ),
        SlideSpec(
            title="A short history of flight (3 moments)",
            bullets=[
                "Balloons: first human flight (buoyancy)",
                "Wright era: powered + controlled flight",
                "Turbines/jets: speed, altitude, long-distance travel",
            ],
            notes="Story-driven, not date-driven. Ask: what enabled long-distance travel?",
        ),
        SlideSpec(
            title="Fixed wing vs rotary wing",
            bullets=[
                "Fixed wing: needs forward motion to make lift",
                "Efficient for long distance",
                "Rotary wing: rotors can hover",
                "Great for point-to-point access",
            ],
            notes="Common misconception: helicopters still generate lift; it’s just from rotor blades.",
        ),
        SlideSpec(
            title="Engines: big picture",
            bullets=[
                "Piston + prop: like a car engine turning a propeller",
                "Turbine: hot gases spin turbine blades",
                "Jet (turbofan): pushes air back to make thrust",
            ],
            notes="Keep it conceptual; avoid deep thermodynamics.",
        ),
        SlideSpec(
            title="The 4 forces of flight",
            bullets=[
                "Lift (up)",
                "Weight (down)",
                "Thrust (forward)",
                "Drag (back)",
            ],
            notes="Draw the arrows on a simple airplane sketch. In level cruise: lift≈weight, thrust≈drag.",
        ),
        SlideSpec(
            title="Lift: how wings make it",
            bullets=[
                "Air moves around an airfoil",
                "Pressure differences help create lift",
                "Wings also push air downward (downwash)",
            ],
            notes="Use paper-strip demo. Phrase it: pressure + downwash.",
        ),
        SlideSpec(
            title="Angle of attack (AOA) + stall",
            bullets=[
                "AOA = angle between wing and airflow",
                "Too high AOA → stall (lift drops, drag rises)",
                "Rule of thumb: many wings stall around ~15–20° AOA",
            ],
            notes=(
                "Clarify: 15–20° is a common rule-of-thumb, not universal; depends on airfoil, flaps, contamination, loading.\n"
                "Key point: stalls are caused by AOA, not speed alone."
            ),
        ),
        SlideSpec(
            title="Control surfaces",
            bullets=[
                "Ailerons → roll",
                "Elevator → pitch",
                "Rudder → yaw",
            ],
            notes="Use hand/airplane model to show roll/pitch/yaw. Mention coordinated turns (bank + appropriate rudder).",
        ),
        SlideSpec(
            title="Cockpit instruments (single-engine)",
            bullets=[
                "Airspeed • Altimeter • Attitude",
                "Heading • Turn/Bank • Vertical speed",
                "Compass • Navigation • Communication",
                "Engine indicators (RPM, temps, pressures, fuel)",
            ],
            notes="Explain what each tells the pilot; emphasize cross-checking multiple instruments.",
        ),
        SlideSpec(
            title="Wrap-up + next steps",
            bullets=[
                "Session 2: safety-first flight principles in action",
                "Drone option: complete FAA TRUST before next session",
                "Bring questions and your handouts",
            ],
            notes="Mention follow-up airport/tower/ARTCC visits if you’re offering them.",
        ),
    ]

    for spec in slides:
        _add_title_and_bullets(prs, spec)

    return prs


def build_session_2() -> Presentation:
    prs = Presentation()

    _add_title_slide(
        prs,
        title="Aviation Merit Badge – Session 2",
        subtitle="Safety briefing • How to fly • Walk out • Drone flight • Debrief",
    )

    slides: List[SlideSpec] = [
        SlideSpec(
            title="Pre-work (do before today)",
            bullets=[
                "Complete FAA TRUST (required to fly)",
                "Fill out preflight sections on the flight log",
                "Wear closed-toe shoes",
            ],
            notes="If someone doesn’t have TRUST, they can observe/record and schedule follow-up flight.",
        ),
        SlideSpec(
            title="Today’s goals",
            bullets=[
                "Fly safely (one pilot at a time)",
                "Successful takeoff → controlled hover → landing",
                "Connect what you see to the 4 forces",
                "Quick overview of aviation pathways (afterward / at home)",
            ],
            notes="Keep content brief; put most instruction before going outside.",
        ),
        SlideSpec(
            title="Safety rules (non-negotiable)",
            bullets=[
                "Visual line of sight (VLOS)",
                "No flying over people",
                "Stay inside the flight box",
                "No stunts—slow and smooth",
                "If anything feels unsafe: LAND",
            ],
            notes="Define the boundary and observer line. Use consistent callouts: taking off, landing, abort, clear.",
        ),
        SlideSpec(
            title="Roles (rotate each flight)",
            bullets=[
                "Pilot",
                "Visual observer",
                "Safety marshal",
                "Recorder",
            ],
            notes="Rotate roles every pilot so everyone participates even with 1 drone.",
        ),
        SlideSpec(
            title="How the controls move the drone",
            bullets=[
                "Throttle: up/down power",
                "Yaw: turn left/right",
                "Pitch/Roll: tilt → move",
            ],
            notes=(
                "Key teaching: tilt trades vertical lift for horizontal motion; that’s why altitude can change when you move.\n"
                "Coaching: small inputs, pause, then adjust."
            ),
        ),
        SlideSpec(
            title="What counts as success today",
            bullets=[
                "Safe takeoff",
                "Stable hover (3–5 seconds)",
                "Safe landing in the box",
            ],
            notes="If time allows, add a slow square pattern. If a scout struggles, shorten the attempt and prioritize safe landing.",
        ),
        SlideSpec(
            title="Walk out + set the flight box",
            bullets=[
                "Walk as a group",
                "Mark boundaries and observer line",
                "One group preflight check",
            ],
            notes="Budget ~10 minutes total for walk + boundary setup (adjust to your site).",
        ),
        SlideSpec(
            title="During flight: quick science callouts",
            bullets=[
                "More speed → more drag",
                "Tilt forward → some lift becomes horizontal motion",
                "Wind → drift and extra corrections",
            ],
            notes="Keep this as short callouts, not a lecture.",
        ),
        SlideSpec(
            title="Debrief questions",
            bullets=[
                "What safety rule mattered most?",
                "What helped stability?",
                "How is a drone like a helicopter?",
            ],
            notes="Use these while walking back or during sign-off time.",
        ),
        SlideSpec(
            title="After class: Requirement 5 (at home)",
            bullets=[
                "Pilot certificates ladder (student → private → instrument → commercial → ATP)",
                "Remote pilot certificate vs TRUST",
                "Choose 1 aviation career to research",
            ],
            notes="Point them to the Req 5 handout for deeper work during the week.",
        ),
        SlideSpec(
            title="Follow-up opportunities",
            bullets=[
                "Airport visit",
                "Control tower visit (if available)",
                "ZKC facility visit (if available)",
            ],
            notes="Mention scheduling constraints and behavior/ID/photo rules.",
        ),
    ]

    for spec in slides:
        _add_title_and_bullets(prs, spec)

    return prs


def build_session_2_requirement_5() -> Presentation:
    """Build an in-class, non-drone deck for Requirement 5.

    Source of truth: instructor guide markdown in instructor-notes.
    Bullets are intentionally concise; the relevant instructor-guide sections are
    attached as speaker notes so you can teach directly from that document.
    """

    prs = Presentation()

    source_md = _read_text(SESSION_2_REQ_5_SOURCE) if SESSION_2_REQ_5_SOURCE.exists() else ""

    _add_title_slide(
        prs,
        title="Aviation Merit Badge – Session 2 (Requirement 5)",
        subtitle="Personal & Professional Aviation Opportunities • classroom module",
    )

    slides: List[SlideSpec] = [
        SlideSpec(
            title="Today’s outcomes",
            bullets=[
                "Explain key certificates/ratings in plain language",
                "Name at least one local youth aviation organization",
                "Identify 3 aviation careers; start researching 1",
            ],
            notes=_notes_from_sections(source_md, ["Outcomes"]),
        ),
        SlideSpec(
            title="Ground rules",
            bullets=[
                "Plain language (no legal fine print)",
                "Rules change — verify current FAA/AOPA details",
                "This is an overview, not flight training",
            ],
            notes=_notes_from_sections(source_md, ["Ground rules (keep it accurate and scout-friendly)"]),
        ),
        SlideSpec(
            title="Choose the version (time check)",
            bullets=[
                "10 min: TRUST vs Part 107 + 1 pilot pathway",
                "30 min: cover all + start careers research",
                "45 min: add activities + quick verbal check-offs",
            ],
            notes=_notes_from_sections(source_md, ["Choose the version that fits your time", "Suggested pacing"]),
        ),
        SlideSpec(
            title="Credential roadmap (big picture)",
            bullets=[
                "Airplane pilot path: Student → Private → Instrument → Commercial → ATP",
                "Instructor path: add CFI",
                "Drone path: TRUST (recreational) vs Part 107 Remote Pilot (work)",
            ],
            notes=_notes_from_sections(source_md, ["Activity A – Pathway ladder (3–5 min)"]),
        ),
        SlideSpec(
            title="TRUST vs Remote Pilot (Part 107)",
            bullets=[
                "TRUST: recreational safety knowledge (required for many hobby flights)",
                "Part 107: common certificate for non-recreational drone operations",
                "Key idea: safety + judgment apply either way",
            ],
            notes=_notes_from_sections(source_md, ["Remote pilot certificate (drone / Part 107)"]),
        ),
        SlideSpec(
            title="Pilot certificates (Part 61) — what they allow",
            bullets=[
                "Student: train with an instructor; first solo with sign-off",
                "Recreational: less common; more limits than private",
                "Private: most common personal flying certificate",
            ],
            notes=_notes_from_sections(
                source_md,
                [
                    "0:03–0:15 – Pilot certificates (Part 61) and what they allow",
                ],
            ),
        ),
        SlideSpec(
            title="Instrument rating — why it matters",
            bullets=[
                "Fly using instruments + procedures when you can’t rely on outside references",
                "Improves discipline, workload management, and weather decision-making",
                "Instrument is a rating (not a separate certificate)",
            ],
            notes=_notes_from_sections(source_md, ["0:15–0:20 – Instrument rating (why it matters)"]),
        ),
        SlideSpec(
            title="Professional certificates",
            bullets=[
                "Commercial: allows certain paid flying jobs (with more requirements)",
                "ATP: top rung for airline captain and many airline roles",
                "CFI: authorized to teach + sign training logs",
            ],
            notes=_notes_from_sections(source_md, ["0:20–0:28 – Professional certificates"]),
        ),
        SlideSpec(
            title="Youth opportunities (local)",
            bullets=[
                "Aviation Exploring Post: learn skills, visit facilities, meet mentors",
                "Civil Air Patrol (CAP): aerospace education + leadership",
                "Action: write 1 org to check + what you hope to do",
            ],
            notes=_notes_from_sections(source_md, ["0:28–0:33 – Youth organizations (local opportunities)"]),
        ),
        SlideSpec(
            title="Careers: choose 3, research 1",
            bullets=[
                "List 3 careers (pilot, mechanic, controller, engineer, dispatcher, airport ops, drone operator)",
                "Pick 1 and research: training, certs, experience, costs",
                "Also research: job outlook, salary range, advancement",
            ],
            notes=_notes_from_sections(source_md, ["0:33–0:45 – Careers: choose 3, research 1"]),
        ),
        SlideSpec(
            title="Quick activity: match credential to scenario",
            bullets=[
                "Learn + fly solo someday → Student",
                "Fly family for fun → Private",
                "Fly in clouds using procedures → Instrument rating",
                "Fly a drone for a business → Part 107 Remote Pilot",
            ],
            notes=_notes_from_sections(source_md, ["Activity B – Match the credential to the scenario (3–6 min)"]),
        ),
        SlideSpec(
            title="Quick check-offs (fast sign-off)",
            bullets=[
                "Each scout answers 1–2 items in one sentence",
                "Rotate prompts so it stays quick",
                "Use the handout as the record",
            ],
            notes=_notes_from_sections(source_md, ["Activity C – One-sentence check-offs (5–10 min)", "Counselor sign-off guidance (simple)"]),
        ),
        SlideSpec(
            title="Wrap-up + next steps",
            bullets=[
                "Finish the Requirement 5 handout (if not done)",
                "Bring 2–3 questions for counselor discussion",
                "Optional follow-ups: airport/tower/facility visit",
            ],
            notes=_notes_from_sections(source_md, ["Counselor sign-off guidance (simple)"]),
        ),
    ]

    for spec in slides:
        _add_title_and_bullets(prs, spec)

    return prs


def main() -> None:
    SLIDES_DIR.mkdir(parents=True, exist_ok=True)

    session_1 = build_session_1()
    session_1_path = SLIDES_DIR / "session-1.pptx"
    session_1.save(session_1_path.as_posix())

    session_2 = build_session_2()
    session_2_path = SLIDES_DIR / "session-2.pptx"
    session_2.save(session_2_path.as_posix())

    session_2_req_5 = build_session_2_requirement_5()
    session_2_req_5_path = SLIDES_DIR / "session-2-req-5.pptx"
    session_2_req_5.save(session_2_req_5_path.as_posix())

    print(f"Wrote {session_1_path.relative_to(ROOT)}")
    print(f"Wrote {session_2_path.relative_to(ROOT)}")
    print(f"Wrote {session_2_req_5_path.relative_to(ROOT)}")


if __name__ == "__main__":
    main()
